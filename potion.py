import asyncio
import json
import logging
import os
from copy import deepcopy
from typing import Callable, List

from fastapi import HTTPException
from pptx import Presentation
from pptx.enum.dml import MSO_FILL
from pptx.slide import Slide
from pydantic import BaseModel, Field, TypeAdapter, ValidationError

from llms.openai import call_openai
from llms.ppt_prompt import design_prompt
from tools.validator import convert_to_filename, trim_code_block

# Configure logging
logging.basicConfig(level=logging.INFO)


class DesignSchema(BaseModel):
    """
    This BaseModel is used for each slide during the design phase.
    """

    slide_title: str = Field(..., description="Name of the slide, e.g., 'Flowchart'")
    desc: str = Field(
        ...,
        description="Description of when to use this slide, e.g., 'To create a flowchart representing time series events'",
    )
    compose_schema_name: str = Field(
        ..., description="Name of the compose schema to be used"
    )


class ComposeSchema(BaseModel):
    """
    This BaseModel is used for each slide during the compose phase.
    """

    name: str = Field(..., description="Name of the slide, e.g., 'Flowchart'")
    desc: str = Field(
        ...,
        description="When to use this slide, e.g., 'To create a flowchart representing time series events'",
    )
    func: Callable = Field(
        ...,
        description="Function to input the schema and create the slide; returns a pptx Slide object",
    )
    slide_layout_index: int = Field(
        ..., description="Which slide layout to use for each slide type"
    )


class Potion:
    def __init__(
        self, template_path: str, compose_schemas: List[ComposeSchema]
    ) -> None:
        self.template_path: str = template_path
        self.compose_schemas: List[ComposeSchema] = compose_schemas

        self.presentation = Presentation(template_path)
        self.ppt_length: int = len(self.presentation.slides)
        logging.info(
            f"Loaded presentation with {self.ppt_length} slides.\nWith ComposeSchemas length of {len(compose_schemas)}"
        )

    def _create_compose_schema_desc(self) -> str:
        """Create a string interpretation of all the compose schemas."""
        res = "**************\n"
        for schema in self.compose_schemas:
            res += f"compose_schema_name: {schema.name}\n"
            res += f"when to call it: {schema.desc}\n"
            res += "**************\n"
        return res

    def get_compose_schema(self, name: str) -> ComposeSchema:
        """Retrieve a ComposeSchema by name."""
        for cs in self.compose_schemas:
            if cs.name == name:
                return cs
        raise HTTPException(500, f"Compose schema '{name}' not found.")

    # async def plan(
    #     self, topic: str
    # ) -> str:
    #     res = await call_openai(f'Generate a title for a powerpoint in the following format on the topic of {topic} do not include " or quotes blocks.')
    #     return res

    async def design(self, query: str, attempts: int = 2) -> List[DesignSchema]:
        """
        From the user description, create a JSON output of the PowerPoint.

        Attributes:
            query (str): user description of the powerpoint, starting from slide 0
                each slide should have a slide title and a description of what the slide will be about.
                the content can be unstructured.
            attempt (int): the number of maximum attempts to retry upon failure to generate / parse

        Return:
            List[DesignSchema]: a list of designschema for composing a powerpoint
        """
        for attempt in range(attempts + 1):
            try:
                # Call to OpenAI API to create the Design Schema
                json_str = await call_openai(
                    design_prompt.format(
                        ppt_outline=query,
                        compose_schema=self._create_compose_schema_desc(),
                    )
                )
                json_str = trim_code_block(json_str)
                logging.info(f"Created JSON outline:\n{json_str}")

                # Load the output into List[DesignSchema] format and return it
                json_outline = json.loads(json_str)
                adapter = TypeAdapter(List[DesignSchema])
                design_schemas = adapter.validate_python(json_outline)
                logging.info("Successfully validated design schemas.")
                return design_schemas

            except (json.JSONDecodeError, ValidationError) as e:
                if attempt == attempts:
                    # Return error details after all attempts fail
                    error_type = (
                        "Invalid JSON format"
                        if isinstance(e, json.JSONDecodeError)
                        else "Validation error"
                    )
                    logging.error(f"{error_type}: {e}")
                    raise HTTPException(500, f"{error_type}: {str(e)}")
                else:
                    logging.warning(f"Attempt {attempt + 1} failed: {e}")
                    continue

    def duplicate_slide(self, slide_index: int) -> Slide:
        """
        Duplicate a slide from the presentation, including all shapes and formatting,
        and remove empty placeholders from the duplicated slide.

        Args:
            slide_index (int): The index of the slide to duplicate.

        Returns:
            Slide: The new slide that is a duplicate of the original one.
        """

        source_slide = self.presentation.slides[slide_index]
        slide_layout = (
            source_slide.slide_layout
        )  # Use the same layout as the source slide
        new_slide = self.presentation.slides.add_slide(slide_layout)

        # Copy shapes from the source slide to the new slide
        for shape in source_slide.shapes:
            new_shape_element = deepcopy(shape.element)
            new_slide.shapes._spTree.insert_element_before(
                new_shape_element, "p:extLst"
            )

        # Remove empty placeholders from the new slide
        self._remove_empty_placeholders(new_slide)

        # Copy slide background (if any)
        self._copy_slide_background(source_slide, new_slide)

        # Copy slide notes (if any)
        if source_slide.has_notes_slide:
            source_notes_slide = source_slide.notes_slide
            new_notes_slide = new_slide.notes_slide
            new_notes_slide.notes_text_frame.text = (
                source_notes_slide.notes_text_frame.text
            )

        logging.info(f"Duplicated slide {slide_index}.")
        return new_slide

    def _copy_slide_background(self, source_slide: Slide, target_slide: Slide) -> None:
        """
        Copy the background from one slide to another.

        Args:
            source_slide (Slide): The slide to copy the background from.
            target_slide (Slide): The slide to copy the background to.
        """
        source_fill = source_slide.background.fill
        target_fill = target_slide.background.fill

        if source_fill.type == MSO_FILL.SOLID:
            target_fill.solid()
            target_fill.fore_color.rgb = source_fill.fore_color.rgb
        elif source_fill.type == MSO_FILL.GRADIENT:
            target_fill.gradient()
            # Copy gradient stops if needed
        elif source_fill.type == MSO_FILL.PICTURE:
            image = source_fill.picture.image
            target_fill.user_picture(image.blob)
        elif source_fill.type == MSO_FILL.PATTERNED:
            target_fill.patterned()
            # Copy pattern details
        elif source_fill.type == MSO_FILL.TEXTURED:
            target_fill.texture()
            # Copy texture details
        else:
            target_fill.background()  # Sets the fill to no fill

    def _remove_empty_placeholders(self, slide: Slide) -> None:
        """
        Remove empty placeholders from a slide.

        Args:
            slide (Slide): The slide from which to remove empty placeholders.
        """
        shapes_to_remove = []
        for shape in slide.shapes:
            if shape.is_placeholder:
                if shape.has_text_frame:
                    if not shape.text.strip():
                        shapes_to_remove.append(shape)
                else:
                    # For non-text placeholders (like image placeholders), check if they are empty
                    shapes_to_remove.append(shape)

        # Remove the collected shapes
        for shape in shapes_to_remove:
            slide.shapes._spTree.remove(shape.element)

    async def compose(self, design_schemas: List[DesignSchema]) -> None:
        """
        Create the PPTX based on the design schema list.
        1. loop through the design schema
            a. get the corresponding compose schema referenced in the design schema
            b. duplicate the slide the compose schema is referencing to to the end.
            c. taking the output from the design schema, feed it into the compose schema function
            d. each function edits the slide inplace.
        2. after all editing is finished. Remove all template slides in the beginning.

        Attributes:
            design_schemas (List[DesignSchema]): a list of designs for composing the slide
                object with.
        """
        tasks = []
        for ds in design_schemas:
            # Find the corresponding ComposeSchema
            compose_schema = self.get_compose_schema(ds.compose_schema_name)

            # Duplicate target slide
            slide = self.duplicate_slide(compose_schema.slide_layout_index)

            # Append the async function to tasks
            tasks.append(compose_schema.func(ds, slide))

        await asyncio.gather(*tasks)

        # Remove the template slides
        self._remove_template_slides()
        logging.info("Completed composing the presentation.")

    def _remove_template_slides(self) -> None:
        """Removes the initial template slides from the presentation."""
        slide_ids = self.presentation.slides._sldIdLst
        for _ in range(self.ppt_length):
            try:
                rId = slide_ids[0].rId
                self.presentation.part.drop_rel(rId)
                del slide_ids[0]
            except Exception as e:
                logging.error(f"Error removing slide: {e}")

    def save(self, filename: str = "output") -> os.PathLike:
        """
        Save the presentation to a file.

        Attribute:
            filename (str): the name of the output file.

        Returns:
            os.PathLike: The path to the saved pptx presentation.
        """

        filename = convert_to_filename(filename)
        output_folder = os.getenv("PPT_OUTPUT_FOLDER", "output/")
        os.makedirs(output_folder, exist_ok=True)

        output_path = os.path.join(output_folder, f"{filename}.pptx")
        self.presentation.save(output_path)
        logging.info(f"Presentation saved to {output_path}.")

        return output_path
