import argparse
import os
import sys
from typing import Dict

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Length

# Custom mapping for placeholder types (based on PowerPoint's common placeholders)
PLACEHOLDER_TYPES: Dict[int, str] = {
    0: "Title",
    1: "Body",
    2: "Center Title",
    3: "Subtitle",
    4: "Date",
    5: "Slide Number",
    6: "Footer",
}


def get_placeholder_type_name(placeholder_type: int) -> str:
    """Return the placeholder type name based on its numeric value."""
    return PLACEHOLDER_TYPES.get(placeholder_type, f"Unknown ({placeholder_type})")


def process_presentation(presentation_path: str) -> None:
    """
    Process and print details of a PowerPoint presentation.

    Args:
        presentation_path (str): Path to the PowerPoint (.pptx) file.
    """
    # Load the presentation
    presentation = Presentation(presentation_path)

    # Iterate over each slide
    for slide_number, slide in enumerate(presentation.slides, start=1):
        print(f"\n========== Slide {slide_number} ==========")

        # Iterate over each shape in the slide
        for shape_number, shape in enumerate(slide.shapes, start=1):
            print(f"\nShape {shape_number}: {shape.name}")
            print(f"- Shape Type: {shape.shape_type}")
            print(
                f"- Position: (Left: {Length(shape.left)}, Top: {Length(shape.top)}, "
                f"Width: {Length(shape.width)}, Height: {Length(shape.height)})"
            )

            # Check if the shape is a placeholder
            if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
                placeholder_type = shape.placeholder_format.type
                placeholder_name = get_placeholder_type_name(placeholder_type)
                print(f"- Placeholder Type: {placeholder_name}")

            # Check if the shape contains a text frame
            if shape.has_text_frame:
                print("- Contains Text:")
                for paragraph in shape.text_frame.paragraphs:
                    print(f"  - Paragraph Text: '{paragraph.text}'")
                    print(f"    Font Size: {paragraph.font.size}")
                    print(f"    Bold: {paragraph.font.bold}")
                    print(f"    Italic: {paragraph.font.italic}")
                    print(f"    Alignment: {paragraph.alignment}")

            # Check if the shape is a picture
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                print("- This is a Picture.")
                print(f"  Image Size: {Length(shape.width)}x{Length(shape.height)}")

            # Check if the shape is an AutoShape or group
            if shape.shape_type in [MSO_SHAPE_TYPE.AUTO_SHAPE, MSO_SHAPE_TYPE.GROUP]:
                print(f"- AutoShape Type: {shape.auto_shape_type}")

        print("\n=========================================\n")


def main() -> None:
    """Parse command-line arguments and process the PowerPoint presentation."""
    parser = argparse.ArgumentParser(description="Process a PowerPoint presentation.")
    parser.add_argument("presentation_path", type=str, help="Path to the .pptx file")

    args = parser.parse_args()

    if not args.presentation_path.endswith(".pptx"):
        print("Error: File must be a .pptx type.")
        sys.exit(1)

    if not os.path.isfile(args.presentation_path):
        print("Error: File does not exist.")
        sys.exit(1)

    process_presentation(args.presentation_path)


if __name__ == "__main__":
    main()
