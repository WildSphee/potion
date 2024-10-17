import markdown2
from bs4 import BeautifulSoup
from pptx.dml.color import RGBColor


def _add_formatted_run(paragraph, text, bold=False, italic=False, color=None) -> None:
    """
    Adds a run with specific formatting to an existing paragraph.
    """
    run = paragraph.add_run()
    run.text = text

    # Apply formatting
    font = run.font
    font.bold = bold  # Bold if needed
    font.italic = italic  # Italic if needed
    if color:
        font.color.rgb = color


def process_markdown_to_ppt(shape, text: str):
    """
    Processes Markdown and adds it to a shape's text frame as a single paragraph.
    """
    # Clear the shape's text content
    shape.text = ""
    text_frame = shape.text_frame

    # Ensure the text frame contains at least one paragraph
    paragraph = text_frame.paragraphs[0]  # Use the first paragraph

    # Convert Markdown to HTML
    html = markdown2.markdown(text)

    # Parse the HTML with BeautifulSoup
    soup = BeautifulSoup(html, "html.parser")

    # Process each element in the parsed HTML
    for element in soup.descendants:
        if element.name == "h1":
            _add_formatted_run(
                paragraph, element.get_text(), bold=True, color=RGBColor(0, 51, 102)
            )
        elif element.name == "h2":
            _add_formatted_run(paragraph, element.get_text(), bold=True)
        elif element.name == "li":
            _add_formatted_run(paragraph, f"- {element.get_text()} ")
        elif element.name == "em":
            _add_formatted_run(paragraph, element.get_text(), italic=True)
        elif element.name == "strong":
            _add_formatted_run(paragraph, element.get_text(), bold=True)
        elif element.name == "p":
            _add_formatted_run(paragraph, element.get_text())

    # Ensure the content fits in the text box (optional)
    text_frame.word_wrap = True

    # add a new line after each -
    shape.text = shape.text.replace("-", "\n-")
